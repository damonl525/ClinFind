import docx
import openpyxl
from pptx import Presentation
from .base import BaseParser
from typing import Tuple

class DocxParser(BaseParser):
    def parse(self, file_path: str) -> Tuple[str, str]:
        doc = docx.Document(file_path)
        full_text = []
        
        # Paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text.strip())
        
        # Tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        full_text.append(cell.text.strip())
                        
        return "\n".join(full_text), ""

class XlsxParser(BaseParser):
    def parse(self, file_path: str) -> Tuple[str, str]:
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        full_text = []
        
        for sheet in wb.worksheets:
            # full_text.append(f"Sheet: {sheet.title}")
            for r_idx, row in enumerate(sheet.iter_rows(values_only=True), 1):
                # Filter None values and convert to string
                row_values = []
                for c_idx, cell in enumerate(row, 1):
                    if cell is not None:
                        # Convert column index to letter (1->A, 2->B...)
                        col_letter = openpyxl.utils.get_column_letter(c_idx)
                        # Format: [Sheet:Name Row:1 Col:A] Value
                        # This ensures the metadata is close to the content for snippet extraction
                        row_values.append(f"[Sheet:{sheet.title} Row:{r_idx} Col:{col_letter}] {cell}")
                
                if row_values:
                    full_text.append(" ".join(row_values))
                    
        return "\n".join(full_text), ""

class PptxParser(BaseParser):
    def parse(self, file_path: str) -> Tuple[str, str]:
        prs = Presentation(file_path)
        full_text = []
        
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            
            # 遍历幻灯片中的所有形状
            for shape in slide.shapes:
                # 处理文本框
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
                
                # 处理表格
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                slide_texts.append(cell.text.strip())
                
                # 处理文本框架（text_frame）
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text:
                            slide_texts.append(para_text)
            
            # 如果这张幻灯片有内容，添加Slide标记
            if slide_texts:
                # 将所有文本合并，并在开头添加Slide标记
                slide_content = " ".join(slide_texts)
                full_text.append(f"[Slide:{i}] {slide_content}")
            else:
                # 即使是空幻灯片，也添加一个标记（可选）
                full_text.append(f"[Slide:{i}] (空白幻灯片)")
                        
        return "\n".join(full_text), ""
